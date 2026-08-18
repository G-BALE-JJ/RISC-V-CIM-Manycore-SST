#!/usr/bin/env python3

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_PATH = OUT_DIR / "final_presentation_cn.pptx"

W = 13.333
H = 7.5

BG = "F6F7F5"
INK = "1D252C"
MUTED = "637078"
GRID = "D8DDDA"
TEAL = "087E8B"
BLUE = "4C78A8"
GREEN = "3D8B66"
RED = "C6423E"
AMBER = "D99A22"
PALE_TEAL = "DDEFF0"
PALE_BLUE = "E4ECF5"
PALE_GREEN = "E2F0E9"
PALE_RED = "F5E3E1"
PALE_AMBER = "F7ECD1"
WHITE = "FFFFFF"

CN_FONT = "Noto Sans CJK SC"
EN_FONT = "Aptos"


def rgb(value):
    return RGBColor.from_string(value)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, x, y, w, h, text, size=16, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=CN_FONT,
             margin=0.04, linesp=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = linesp
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, kicker=None, num=None):
    if kicker:
        add_text(slide, 0.55, 0.22, 7.0, 0.28, kicker.upper(), 8.5, TEAL, True, font=EN_FONT)
    add_text(slide, 0.55, 0.48, 11.8, 0.58, title, 26, INK, True)
    if num is not None:
        add_text(slide, 12.30, 0.38, 0.45, 0.35, f"{num:02d}", 9, MUTED, True,
                 PP_ALIGN.RIGHT, font=EN_FONT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.10), Inches(12.2), Inches(0.018))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(GRID); line.line.fill.background()


def rect(slide, x, y, w, h, fill, line=GRID, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def label_box(slide, x, y, w, h, title, body="", fill=WHITE, accent=TEAL,
              title_size=15, body_size=10.5):
    rect(slide, x, y, w, h, fill, GRID, True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(accent); bar.line.fill.background()
    add_text(slide, x + 0.18, y + 0.12, w - 0.32, 0.34, title, title_size, INK, True)
    if body:
        add_text(slide, x + 0.18, y + 0.52, w - 0.32, h - 0.62, body, body_size, MUTED)


def arrow(slide, x, y, w, h=0.28, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = rgb(color); a.line.fill.background()
    return a


def metric(slide, x, y, w, value, label, color=TEAL):
    add_text(slide, x, y, w, 0.48, value, 25, color, True, font=EN_FONT)
    add_text(slide, x, y + 0.49, w, 0.36, label, 10.5, MUTED)


def takeaway(slide, text, color=INK, fill="E9ECE9"):
    rect(slide, 0.55, 6.88, 12.2, 0.40, fill, fill, True)
    add_text(slide, 0.72, 6.94, 11.85, 0.25, text, 11.5, color, True,
             valign=MSO_ANCHOR.MIDDLE)


def source(slide, text):
    add_text(slide, 0.58, 7.30, 11.9, 0.14, text, 7.2, MUTED, font=EN_FONT)


def notes(slide, text):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text


def bar(slide, x, y, w, h, value, max_value, color, label, value_text):
    rect(slide, x, y, w, h, "E5E8E6", "E5E8E6", True)
    fill_w = max(0.04, w * value / max_value)
    rect(slide, x, y, fill_w, h, color, color, True)
    add_text(slide, x, y - 0.30, w * 0.58, 0.24, label, 10.5, INK, True)
    add_text(slide, x + w * 0.58, y - 0.30, w * 0.42, 0.24, value_text, 10.5, color, True,
             PP_ALIGN.RIGHT, font=EN_FONT)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


# 1. Cover
slide = prs.slides.add_slide(blank); set_bg(slide)
add_text(slide, 0.65, 0.55, 7.8, 0.35, "GROUP MEETING · 2026.07.21", 10, TEAL, True, font=EN_FONT)
add_text(slide, 0.65, 1.18, 7.6, 1.35, "面向 NoC Manycore 的\nSoftmax Row Engine", 35, INK, True)
add_text(slide, 0.68, 2.72, 6.8, 0.64, "从显式 NoC reduction 到单 tensor job、\n行级并行与四上下文流水", 17, MUTED)
add_text(slide, 0.68, 6.44, 5.8, 0.34, "RISC-V · Vanadis · SST · Golem", 11, INK, True, font=EN_FONT)

# Cover visual: job -> mesh -> HBM
rect(slide, 8.35, 0.72, 4.22, 5.85, WHITE, GRID, True)
label_box(slide, 9.05, 1.10, 2.82, 0.72, "1 Tensor Job", "hardware row scheduling", PALE_TEAL, TEAL, 15, 9)
arrow(slide, 10.22, 1.93, 0.52, 0.24, TEAL)
for r in range(4):
    for c in range(4):
        x = 8.82 + c * 0.78; y = 2.42 + r * 0.70
        color = [PALE_BLUE, PALE_TEAL, PALE_GREEN, PALE_AMBER][c]
        rect(slide, x, y, 0.58, 0.49, color, GRID, True)
        add_text(slide, x, y + 0.12, 0.58, 0.18, f"T{r*4+c}", 8, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
for c, col in enumerate([BLUE, TEAL, GREEN, AMBER]):
    x = 8.79 + c * 0.78
    rect(slide, x, 5.52, 0.64, 0.48, col, col, True)
    add_text(slide, x, 5.65, 0.64, 0.16, f"HBM{c}", 8, WHITE, True, PP_ALIGN.CENTER, font=EN_FONT)
add_text(slide, 8.86, 6.13, 3.3, 0.24, "16 Row Engines · 4 HBM Nodes", 9.5, MUTED, True, PP_ALIGN.CENTER, font=EN_FONT)
notes(slide, "今天汇报的是Softmax加速架构从旧的distributed-column路径，到Row Engine和单tensor job路径的完整进展。重点不是单一cycle数字，而是为什么架构变化能够把控制、reduction和DMA开销系统性消除。")


# 2. Problem
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "旧方案把一行拆给16个worker，通信和控制随规模膨胀", "Problem", 2)
add_text(slide, 0.62, 1.35, 3.3, 0.32, "旧 distributed-column 数据流", 14, INK, True)
rect(slide, 0.62, 1.82, 3.25, 0.62, PALE_BLUE, GRID, True)
add_text(slide, 0.80, 2.02, 2.90, 0.22, "1 row × 4096 columns", 13, BLUE, True, PP_ALIGN.CENTER, font=EN_FONT)
arrow(slide, 2.02, 2.54, 0.48, 0.24, BLUE)
for i in range(8):
    x = 0.63 + (i % 4) * 0.80; y = 3.02 + (i // 4) * 0.60
    rect(slide, x, y, 0.62, 0.42, PALE_BLUE, GRID, True)
    add_text(slide, x, y + 0.11, 0.62, 0.16, f"W{i*2}-{i*2+1}", 7.5, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
arrow(slide, 2.02, 4.33, 0.48, 0.24, RED)
label_box(slide, 0.72, 4.82, 3.02, 0.75, "MAX / SUM 两次全局归约", "barrier + request/response", PALE_RED, RED, 13, 9.5)
arrow(slide, 2.02, 5.68, 0.48, 0.24, RED)
add_text(slide, 0.72, 6.05, 3.02, 0.36, "Normalize 后再写回 HBM", 11, INK, True, PP_ALIGN.CENTER)

label_box(slide, 4.42, 1.44, 3.60, 1.15, "错误的并行层级", "同一行上的worker越多，fan-in、barrier和等待越重", WHITE, RED)
label_box(slide, 4.42, 2.88, 3.60, 1.15, "任务与DMA过细", "4,096 jobs；input/output各16,384次DMA", WHITE, AMBER)
label_box(slide, 4.42, 4.32, 3.60, 1.15, "仿真与硬件进度耦合", "wait polling推进状态，Vanadis退休指令显著增加", WHITE, BLUE)

metric(slide, 8.62, 1.55, 3.6, "65,536", "reduction request messages", RED)
metric(slide, 8.62, 2.82, 3.6, "16.64 M", "whole-system cycles @ 2.3 GHz", BLUE)
metric(slide, 8.62, 4.09, 3.6, "41.9 M", "Vanadis retired instructions", AMBER)
takeaway(slide, "核心问题不是 Softmax 数学本身，而是行/列映射、任务粒度和完成机制不匹配。", RED, PALE_RED)
source(slide, "Source: architecture design §2.2; retained 1024×4096 explicit-NoC capacity run")
notes(slide, "旧方案功能正确，但把4096维的一行拆给16个worker，必须进行global max和global sum两次归约。1024行最终产生65536个reduction request、4096个job和大量小DMA。这里要强调：当时NoC利用率并不高，所以根因不是链路不够快，而是并行组织方式错误。")


# 3. Evolution
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "三轮架构收敛将完整系统周期降低27.8×", "Progress", 3)
stages = [
    ("Explicit-NoC baseline", "16.64 M", "16 workers / row\n4,096 jobs\n65,536 reductions", RED),
    ("Row Engine v2", "1.137 M", "1 tile / row\n16 jobs\n4 contexts", BLUE),
    ("Tensor controller", "598 k", "1 tensor job\n16 endpoints\ncompletion aggregation", GREEN),
]
for i, (name, cyc, body, col) in enumerate(stages):
    x = 0.70 + i * 4.18
    rect(slide, x, 1.58, 3.55, 4.35, WHITE, GRID, True)
    add_text(slide, x + 0.22, 1.86, 3.10, 0.35, name, 14, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x + 0.22, 2.45, 3.10, 0.65, cyc, 31, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x + 0.22, 3.28, 3.10, 1.10, body, 12, INK, False, PP_ALIGN.CENTER)
    if i < 2:
        arrow(slide, x + 3.68, 3.33, 0.52, 0.28, TEAL)
add_text(slide, 8.99, 5.17, 2.95, 0.30, "最终 accelerator path", 10, MUTED, True, PP_ALIGN.CENTER)
add_text(slide, 8.99, 5.48, 2.95, 0.42, "66,062 cycles", 20, TEAL, True, PP_ALIGN.CENTER, font=EN_FONT)
takeaway(slide, "系统级改善来自架构收敛；66,062 cycles 是 descriptor 到完成的 accelerator 指标。", INK, "E8EFED")
source(slide, "Whole-system values use comparable retained runs; clean instrumentation build is excluded from this comparison.")
notes(slide, "这页给出主线：初始whole-system约1664万cycle；Row Engine v2降到113.7万；单tensor controller进一步到59.8万，整体约27.8倍。最终硬件任务路径是66062cycle。clean instrumentation版本会增加诊断输出，因此whole-system不放进这张对比图。")


# 4. Overall architecture
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "单 tensor job 在硬件内分发1024行，16个Row Engine独立推进", "Architecture", 4)
label_box(slide, 4.48, 1.35, 4.30, 0.78, "Tensor Softmax Job Controller", "row allocation · credits · completion aggregation", PALE_TEAL, TEAL, 16, 9.5)
arrow(slide, 6.35, 2.24, 0.54, 0.28, TEAL)
for i in range(16):
    r, c = divmod(i, 8)
    x = 0.68 + c * 1.55; y = 2.82 + r * 1.18
    color = [PALE_BLUE, PALE_TEAL, PALE_GREEN, PALE_AMBER][c % 4]
    rect(slide, x, y, 1.22, 0.82, color, GRID, True)
    add_text(slide, x, y + 0.12, 1.22, 0.22, f"Tile {i}", 10, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x, y + 0.43, 1.22, 0.18, "64 rows", 8.5, MUTED, True, PP_ALIGN.CENTER, font=EN_FONT)
arrow(slide, 6.35, 5.13, 0.54, 0.28, TEAL)
for i, (name, col) in enumerate(zip(["HBM0", "HBM1", "HBM2", "HBM3"], [BLUE, TEAL, GREEN, AMBER])):
    x = 3.05 + i * 1.85
    rect(slide, x, 5.66, 1.45, 0.64, col, col, True)
    add_text(slide, x, 5.86, 1.45, 0.20, name, 11, WHITE, True, PP_ALIGN.CENTER, font=EN_FONT)
add_text(slide, 10.58, 5.58, 1.8, 0.54, "band-striped\ninput + output", 9.5, MUTED, True, PP_ALIGN.CENTER, font=EN_FONT)
takeaway(slide, "dim=4096 主路径采用 1 tile/row；每个tile完整拥有一行，因此不产生跨tile reduction。")
source(slide, "Source: row-engine architecture design §§4, 6, 9")
notes(slide, "coordinator只提交一个128-byte descriptor和一个versioned params block。controller按band将1024行分到16个tile，每个tile处理64行。对4096维而言，一行完整放在一个tile内，因此max和sum都在tile内部完成，NoC上不再发送reduction消息。")


# 5. Row engine microarchitecture
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "Row Engine 将vector、reduction、EXP和DMA拆成独立资源", "Microarchitecture", 5)
modules = [
    ("Input DMA", "HBM → SRAM", BLUE, PALE_BLUE),
    ("16-lane Vector", "MAX / SUB / MUL", TEAL, PALE_TEAL),
    ("Reduction Tree", "16 inputs · 4 cycles", GREEN, PALE_GREEN),
    ("4-lane EXP", "latency 8 · pipelined", RED, PALE_RED),
    ("Reciprocal", "1 / row_sum", AMBER, PALE_AMBER),
    ("Output DMA", "SRAM → HBM", BLUE, PALE_BLUE),
]
for i, (title, body, col, pale) in enumerate(modules):
    x = 0.64 + i * 2.08
    rect(slide, x, 2.10, 1.72, 1.24, pale, GRID, True)
    add_text(slide, x + 0.08, 2.33, 1.56, 0.28, title, 12, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x + 0.08, 2.75, 1.56, 0.25, body, 8.5, MUTED, False, PP_ALIGN.CENTER, font=EN_FONT)
    if i < len(modules)-1:
        arrow(slide, x + 1.76, 2.58, 0.31, 0.20, MUTED)
rect(slide, 2.06, 4.05, 9.20, 1.13, WHITE, GRID, True)
add_text(slide, 2.28, 4.26, 2.25, 0.30, "64 KiB Scratchpad", 15, INK, True, font=EN_FONT)
for i in range(4):
    x = 4.82 + i * 1.42
    rect(slide, x, 4.24, 1.14, 0.56, [PALE_BLUE, PALE_TEAL, PALE_GREEN, PALE_AMBER][i], GRID, True)
    add_text(slide, x, 4.42, 1.14, 0.18, f"Context {i}", 8.5, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
add_text(slide, 9.98, 4.31, 1.04, 0.28, "16 KiB / row", 8.5, MUTED, True, PP_ALIGN.CENTER, font=EN_FONT)
label_box(slide, 0.80, 5.65, 3.30, 0.72, "Event-driven timing", "stage completion self-event", WHITE, TEAL, 12.5, 9)
label_box(slide, 4.98, 5.65, 3.30, 0.72, "Resource reservation", "free_cycle + ready_cycle", WHITE, GREEN, 12.5, 9)
label_box(slide, 9.10, 5.65, 3.30, 0.72, "No per-element event", "降低SST simulation time", WHITE, AMBER, 12.5, 9)
takeaway(slide, "资源解耦允许不同row同时处于MAX、EXP、Normalize和DMA阶段。")
source(slide, "Source: row-engine architecture design §5; current SST parameters")
notes(slide, "每个SFU内部不再是一个monolithic softmax操作槽，而是把vector ALU、reduction tree、EXP pipeline、reciprocal和DMA分别保留free cycle。仿真采用事件驱动的stage completion，而不是每元素每周期发一个SST event，这同时改善架构表达和simulation time。")


# 6. Scratchpad dataflow
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "一行只读一次、只写一次，中间结果全部留在16 KiB本地buffer", "Dataflow", 6)
steps = [
    ("1", "DMA IN", "buffer = x", BLUE, PALE_BLUE),
    ("2", "MAX", "x 保留\n产生 local_max", TEAL, PALE_TEAL),
    ("3", "EXP + SUM", "覆盖为 exp(x-max)\n累加 local_sum", RED, PALE_RED),
    ("4", "NORMALIZE", "buffer *= 1/sum", GREEN, PALE_GREEN),
    ("5", "DMA OUT", "写回 HBM", AMBER, PALE_AMBER),
]
for i, (num, title, body, col, pale) in enumerate(steps):
    x = 0.58 + i * 2.52
    rect(slide, x, 2.02, 2.05, 2.45, pale, GRID, True)
    add_text(slide, x + 0.16, 2.20, 0.42, 0.42, num, 20, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x + 0.20, 2.82, 1.65, 0.34, title, 14, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x + 0.20, 3.42, 1.65, 0.62, body, 10, MUTED, False, PP_ALIGN.CENTER)
    if i < 4:
        arrow(slide, x + 2.08, 3.08, 0.38, 0.24, col)
metric(slide, 1.03, 5.22, 2.5, "16 MiB", "input payload", BLUE)
metric(slide, 4.20, 5.22, 2.5, "0", "intermediate HBM traffic", GREEN)
metric(slide, 7.40, 5.22, 2.5, "16 MiB", "output payload", AMBER)
add_text(slide, 10.43, 5.36, 2.1, 0.64, "FP32 1024×4096\n总不可消除流量 32 MiB", 11, INK, True, PP_ALIGN.CENTER)
takeaway(slide, "Softmax仍需三次本地数学pass，但不再为中间EXP结果访问GlobalMemory/HBM。")
source(slide, "Source: architecture design §§3.3, 5.2")
notes(slide, "稳定Softmax需要先求max，再计算exp和sum，再normalize，因此三次数学pass不能简单消失。优化点是把一整行保留在16KiB scratchpad：EXP阶段原地覆盖输入并同时累加sum，normalize直接读取本地EXP结果，最终只保留一次输入和一次输出HBM流量。")


# 7. Context pipeline
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "四个row context把三阶段串行延迟转化为稳态流水吞吐", "Scheduling", 7)
add_text(slide, 0.72, 1.45, 1.30, 0.30, "时间 →", 10, MUTED, True)
colors = {"DMA": BLUE, "MAX": TEAL, "EXP": RED, "NORM": GREEN, "OUT": AMBER}
schedule = [
    [("DMA", 0.0, .7), ("MAX", .8, .8), ("EXP", 1.7, 2.0), ("NORM", 3.8, .8), ("OUT", 4.7, .7)],
    [("DMA", .8, .7), ("MAX", 1.7, .8), ("EXP", 2.6, 2.0), ("NORM", 4.7, .8), ("OUT", 5.6, .7)],
    [("DMA", 1.7, .7), ("MAX", 2.6, .8), ("EXP", 3.5, 2.0), ("NORM", 5.6, .8)],
    [("DMA", 2.6, .7), ("MAX", 3.5, .8), ("EXP", 4.4, 2.0)],
]
scale = 1.55
for row, tasks in enumerate(schedule):
    y = 2.02 + row * 0.88
    add_text(slide, 0.67, y + 0.12, 1.0, 0.22, f"Context {row}", 10, INK, True, font=EN_FONT)
    rect(slide, 1.72, y, 10.35, 0.52, "EAEEEB", "EAEEEB", True)
    for name, start, dur in tasks:
        x = 1.76 + start * scale; w = dur * scale - 0.04
        rect(slide, x, y + 0.04, w, 0.44, colors[name], colors[name], True)
        add_text(slide, x, y + 0.15, w, 0.16, name, 7.8, WHITE, True, PP_ALIGN.CENTER, font=EN_FONT)
for i, (name, col) in enumerate(colors.items()):
    x = 1.88 + i * 1.55
    rect(slide, x, 5.78, 0.22, 0.22, col, col, True)
    add_text(slide, x + 0.30, 5.76, 1.05, 0.24, name, 9, INK, True, font=EN_FONT)
metric(slide, 9.70, 5.50, 2.3, "32.8%", "measured modeled overlap", TEAL)
takeaway(slide, "1024点将98,304 active cycles压缩为66,048-cycle temporal span。")
source(slide, "Schematic is conceptual; overlap percentage comes from the fresh 1024×4096 stage statistics.")
notes(slide, "单行依然需要MAX、EXP和normalize，但多个row可以重叠。四个context正好覆盖不同阶段和DMA等待。新统计显示1024点三个阶段active cycles之和是98304，而时间跨度是66048，重叠节省32256，也就是32.8%。")


# 8. NoC/HBM layout
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "4×4计算mesh按列连接4个HBM节点，连续band保持局部性", "NoC & HBM", 8)
for r in range(4):
    for c in range(4):
        x = 1.00 + c * 1.45; y = 1.58 + r * 1.00
        col = [BLUE, TEAL, GREEN, AMBER][c]
        pale = [PALE_BLUE, PALE_TEAL, PALE_GREEN, PALE_AMBER][c]
        rect(slide, x, y, 0.90, 0.62, pale, col, True)
        add_text(slide, x, y + 0.20, 0.90, 0.18, f"T{r*4+c}", 9, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
        if c < 3:
            arrow(slide, x + 0.94, y + 0.21, 0.42, 0.18, GRID)
for c, (name, col) in enumerate(zip(["HBM0", "HBM1", "HBM2", "HBM3"], [BLUE, TEAL, GREEN, AMBER])):
    x = 0.88 + c * 1.45
    rect(slide, x, 5.72, 1.12, 0.56, col, col, True)
    add_text(slide, x, 5.90, 1.12, 0.18, name, 9, WHITE, True, PP_ALIGN.CENTER, font=EN_FONT)

label_box(slide, 7.15, 1.55, 4.95, 1.02, "Band striping", "band = row / 64\nnode = 1 + band mod 4", WHITE, TEAL, 15, 11)
label_box(slide, 7.15, 2.93, 4.95, 1.02, "连续地址", "每个tile的64行连续存放，支持256 KiB burst", WHITE, BLUE, 15, 11)
label_box(slide, 7.15, 4.31, 4.95, 1.02, "流量不复制", "input/output各16 MiB，四节点分担backend压力", WHITE, GREEN, 15, 11)
takeaway(slide, "HBM striping是带宽约束下的必要结构，而不是单纯增加NoC峰值参数。")
source(slide, "Source: architecture design §8; implementation uses contiguous tile-band layout")
notes(slide, "四个data HBM node与mesh四列对齐。实际实现采用连续tile-band布局：每64行一个band，band轮转到四个节点。这样既保持每个tile访问连续，又让四个backend分担32MiB的不可消除流量。")


# 9. Control and lifecycle
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "完成通知由DMA ACK驱动，guest wait不再推进硬件状态", "Control & ABI", 9)
actors = [("Vanadis guest", BLUE), ("RoCC / SFU", TEAL), ("Job Controller", GREEN), ("16 Row Engines", RED), ("HBM / NoC", AMBER)]
for i, (name, col) in enumerate(actors):
    x = 0.55 + i * 2.53
    add_text(slide, x, 1.42, 2.05, 0.32, name, 11, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+1.00), Inches(1.86), Inches(0.012), Inches(4.30))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(GRID); line.line.fill.background()
events = [
    (0, 1, 2.05, "128-B desc + 64-B params"),
    (1, 2, 2.72, "accept / decode"),
    (2, 3, 3.39, "16 band dispatches"),
    (3, 4, 4.06, "64 input + 64 output DMA"),
    (4, 3, 4.73, "output ACK"),
    (3, 2, 5.30, "16 completion messages"),
    (2, 0, 5.87, "single tensor completion"),
]
for a, b, y, textv in events:
    x1 = 1.55 + a * 2.53; x2 = 1.55 + b * 2.53
    width = x2 - x1
    if width > 0:
        arrow(slide, x1, y, width - 0.10, 0.20, actors[b][1])
    else:
        ar = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(x2+0.10), Inches(y), Inches(-width-0.10), Inches(0.20))
        ar.fill.solid(); ar.fill.fore_color.rgb = rgb(actors[b][1]); ar.line.fill.background()
    add_text(slide, min(x1,x2)+0.10, y-0.27, abs(width)-0.20, 0.22, textv, 8.5, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
takeaway(slide, "descriptor acceptance到最后output ACK为17,565 cycles；完成语义不依赖固定延迟。")
source(slide, "Source: SFUSoftmaxJobParamsV1 ABI and fresh tensor-controller timeline")
notes(slide, "guest只负责构造descriptor和params并发出一次job。controller发送16个band dispatch，每个endpoint内部完成DMA和计算。只有最后一个output DMA ACK到达后，controller才聚合完成并让guest wait返回。这个机制替代了固定store drain和高频poll推进。")


# 10. Test setup
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "测试同时约束功能正确性、生命周期和周期口径", "Evaluation", 10)
label_box(slide, 0.65, 1.52, 3.75, 1.58, "Workload", "FP32 row-wise Softmax\n1024 × 4096\n4,194,304 elements", WHITE, BLUE, 16, 12)
label_box(slide, 4.80, 1.52, 3.75, 1.58, "Architecture", "16 physical SFUs · 4 contexts\n16 vector lanes · 4 EXP lanes\n64 KiB scratchpad", WHITE, TEAL, 16, 11.5)
label_box(slide, 8.95, 1.52, 3.75, 1.58, "Fabric", "4×4 compute mesh · 4 HBM nodes\n1200 GB/s link/xbar\n256 KiB DMA burst · 2.3 GHz", WHITE, GREEN, 16, 11.5)

checks = [
    ("Golden", "4,194,304 checked\n0 mismatch", GREEN),
    ("Lifecycle", "1 job · 1024 rows\n64 DMA pairs", TEAL),
    ("Errors", "0 retry / stale\n0 reduction", BLUE),
    ("Timing", "SST ticks + Vanadis rdcycle\nseparate contracts", AMBER),
]
for i, (title, body, col) in enumerate(checks):
    x = 0.72 + i * 3.05
    rect(slide, x, 3.70, 2.58, 1.55, "EDEFEA", GRID, True)
    add_text(slide, x+0.15, 3.98, 2.28, 0.28, title, 13, col, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x+0.15, 4.45, 2.28, 0.48, body, 10, MUTED, False, PP_ALIGN.CENTER, font=EN_FONT)
metric(slide, 1.15, 5.62, 2.7, "66 tests", "focused unit/contract suite", TEAL)
metric(slide, 5.15, 5.62, 2.7, "4 shapes", "16 / 64 / 256 / 1024 rows", BLUE)
metric(slide, 9.15, 5.62, 2.7, "PASS", "fresh build + diff + parser", GREEN)
takeaway(slide, "任何cycle结果只有在golden、事件计数和completion contract同时通过时才被接受。")
source(slide, "Source: muticore_softmax runner, parser and fresh 2026-07-21 artifacts")
notes(slide, "测试不是只看输出数值。每个点都要求全量logits golden、行数和job数一致、band/worker/completion事件数量匹配、DMA ACK完整、retry和stale为零。周期同时保留SST时序和guest rdcycle，避免把本地wall time或picosecond tick误称为硬件cycle。")


# 11. Key experiments
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "三个负实验排除了“加worker、改burst、去polling”三条伪瓶颈", "Key Experiments", 11)
label_box(slide, 0.65, 1.48, 3.72, 3.90, "A · Workers per row", "16×4096 simulated time\n\n4 workers   422.0 μs\n8 workers   423.4 μs\n16 workers  427.1 μs\n\n更多worker增加fan-in，没有加速。", WHITE, RED, 16, 11.5)
label_box(slide, 4.80, 1.48, 3.72, 3.90, "B · DMA burst", "1024×4096 issue-to-completion\n\n256 KiB  357,385 cycles\n64 KiB    365,093 cycles\n\n64 KiB虽减少NoC stall，却让主指标回退。", WHITE, AMBER, 16, 11.5)
label_box(slide, 8.95, 1.48, 3.72, 3.90, "C · Wait polling", "compatibility polls\n\n1,056,100 → 0\ncycle仍为357,385\n\nPolling影响simulation开销，但不是架构关键路径。", WHITE, BLUE, 16, 11.5)
add_text(slide, 1.02, 5.79, 11.3, 0.34, "证据驱动的结论：必须改变row scheduling和tensor-level control，而不是继续调外围参数。", 15, INK, True, PP_ALIGN.CENTER)
takeaway(slide, "负实验帮助把优化方向从参数DSE收敛到架构重构。", INK, PALE_AMBER)
source(slide, "Source: durable findings; Row Engine v2 controlled experiments")
notes(slide, "这一页适合重点讲。第一，4到16 workers并没有加速16×4096。第二，64KiB burst降低了stall总数，但issue-to-completion反而变差。第三，poll从105万降到0后周期完全不变。这三组负结果共同证明，真正需要的是tensor controller和row-level mapping。")


# 12. Timeline
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "最终路径已由NoC/DMA受限转为Row Engine计算模型受限", "Target Timeline", 12)
axis_x, axis_y, axis_w = 1.05, 2.10, 10.95
rect(slide, axis_x, axis_y, axis_w, 0.06, GRID, GRID)
events_t = [
    (0, "descriptor accepted", TEAL),
    (17565, "final ACK / completion\n17,565 / 17,612", AMBER),
    (66062, "modeled ready", RED),
]
max_t = 66062
for t, lab, col in events_t:
    x = axis_x + axis_w * t / max_t
    rect(slide, x, axis_y-0.13, 0.05, 0.32, col, col)
    align = PP_ALIGN.LEFT if t < 40000 else PP_ALIGN.RIGHT
    tx = x if align == PP_ALIGN.LEFT else x-1.7
    label_text = f"{lab}\n{t:,}" if t not in (17565,) else lab
    add_text(slide, tx, axis_y+0.22, 1.75, 0.48, label_text, 8.5, col, True, align, font=EN_FONT)

bar(slide, 1.05, 3.52, 10.95, 0.40, 17565, 66062, AMBER, "Descriptor → final output DMA ACK", "17,565 cycles")
bar(slide, 1.05, 4.43, 10.95, 0.40, 66062, 66062, RED, "Descriptor → modeled compute ready", "66,062 cycles")
bar(slide, 1.05, 5.34, 10.95, 0.40, 6398, 66062, BLUE, "Clean guest pre-accept setup", "6,398 cycles")
metric(slide, 9.78, 6.02, 2.2, "3.76×", "compute / DMA interval", RED)
takeaway(slide, "所有output ACK在17.6k cycles内完成；66.1k关键路径来自当前4-lane EXP吞吐模型。", RED, PALE_RED)
source(slide, "Fresh 1024×4096 timeline: muticore_softmax_clean_kernel_r1024_d4096_stage")
notes(slide, "descriptor接受后，第一worker dispatch只需11cycle，所有output ACK在17565cycle到达，completion再晚47cycle。但modeled compute ready是66062cycle。因此当前配置下DMA和NoC已经被计算模型隐藏，主瓶颈转到EXP阶段。clean guest在descriptor前还有6398cycle，主要是helper进入和descriptor构造。")


# 13. Scaling
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "固定架构下周期随行数稳定增长，模型与观测只差1 cycle", "Shape Scaling", 13)
rows = [16, 64, 256, 1024]
accel = [1550, 4622, 16910, 66062]
kernel = [7819, 10954, 23323, 72409]
maxv = 72409
base_y = 6.05
chart_x, chart_w, chart_h = 0.95, 8.45, 4.55
for tick in [0, 20000, 40000, 60000]:
    y = base_y - chart_h * tick / maxv
    rect(slide, chart_x, y, chart_w, 0.012, GRID, GRID)
    add_text(slide, 0.55, y-0.12, 0.36, 0.22, f"{tick//1000}k", 8, MUTED, False, PP_ALIGN.RIGHT, font=EN_FONT)
for i, r in enumerate(rows):
    x = 1.38 + i * 2.00
    ah = chart_h * accel[i] / maxv
    kh = chart_h * kernel[i] / maxv
    rect(slide, x, base_y-ah, 0.58, ah, TEAL, TEAL, True)
    rect(slide, x+0.72, base_y-kh, 0.58, kh, BLUE, BLUE, True)
    add_text(slide, x-0.10, 6.20, 1.50, 0.26, str(r), 9.5, INK, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x-0.22, base_y-ah-0.30, 0.82, 0.22, f"{accel[i]:,}", 8, TEAL, True, PP_ALIGN.CENTER, font=EN_FONT)
    add_text(slide, x+0.61, base_y-kh-0.30, 0.82, 0.22, f"{kernel[i]:,}", 8, BLUE, True, PP_ALIGN.CENTER, font=EN_FONT)
rect(slide, 9.95, 2.05, 0.24, 0.24, TEAL, TEAL, True); add_text(slide, 10.29, 2.02, 2.1, 0.26, "Accelerator observed", 10, INK, True, font=EN_FONT)
rect(slide, 9.95, 2.58, 0.24, 0.24, BLUE, BLUE, True); add_text(slide, 10.29, 2.55, 2.1, 0.26, "Clean kernel window", 10, INK, True, font=EN_FONT)
label_box(slide, 9.70, 3.35, 2.62, 1.32, "稳态斜率", "每增加16行\n约 +1,024 cycles", WHITE, RED, 14, 11)
label_box(slide, 9.70, 4.98, 2.62, 1.02, "Contract", "4 shapes 全部 PASS", WHITE, GREEN, 14, 10.5)
takeaway(slide, "Scaling曲线符合四上下文流水模型，没有出现1024点的非线性恶化。")
source(slide, "Fresh fixed-parameter runs: 16/64/256/1024 × 4096, 2026-07-21")
notes(slide, "这里固定16个物理SFU、4 contexts、4个HBM节点、256KiB burst和2.3GHz，只改变行数。accelerator observed分别是1550、4622、16910和66062；与modeled值始终只差1cycle。说明周期模型和实际完成边界一致，1024点没有额外的非线性拖慢。")


# 14. GPU reference
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "与GPU比较只能作为量级参考，必须统一测量边界", "GPU Reference", 14)
vals = [("Golem accelerator", 66062, TEAL), ("Golem clean kernel", 72409, BLUE), ("TileLang p50", 115718, GREEN), ("PyTorch p50", 207911, AMBER)]
maxg = 207911
for i, (name, val, col) in enumerate(vals):
    y = 1.78 + i * 1.08
    bar(slide, 0.95, y, 8.70, 0.48, val, maxg, col, name, f"{val:,} cycles")
label_box(slide, 10.02, 1.52, 2.45, 1.20, "同一shape", "1024 × 4096\nrow-wise Softmax", WHITE, TEAL, 14, 11)
label_box(slide, 10.02, 3.05, 2.45, 1.55, "不同口径", "GPU：nominal SM clock估算\nGolem：SST/rdcycle事件边界", WHITE, RED, 14, 10)
label_box(slide, 10.02, 4.94, 2.45, 1.16, "不宣称", "cycle更少 ≠ 硬件更快", WHITE, AMBER, 14, 10.5)
takeaway(slide, "当前结果说明模型已进入GPU量级；公平比较仍需频率、吞吐、面积和精度共同约束。", INK, PALE_AMBER)
source(slide, "GPU values supplied from 1024×4096 measurements; cycles estimated from nominal SM clock.")
notes(slide, "这页要保守表达。GPU数据来自已有1024×4096测试，TileLang p50约115718 nominal cycles，PyTorch约207911。Golem accelerator和clean kernel分别是66062和72409。但GPU是按标称SM频率从微秒换算，Golem是SST事件和rdcycle窗口，不能据此宣称真实芯片更快。")


# 15. Bottleneck and next architecture
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "下一步应扩展EXP吞吐，而不是继续提高NoC带宽", "Bottleneck & Next", 15)
# Stage intervals
add_text(slide, 0.72, 1.45, 5.4, 0.30, "当前1024×4096 stage active cycles", 13, INK, True)
bar(slide, 0.75, 2.05, 5.35, 0.44, 16384, 65536, TEAL, "MAX · 16 lanes", "16,384")
bar(slide, 0.75, 3.12, 5.35, 0.44, 65536, 65536, RED, "EXP + SUM · 4 lanes", "65,536")
bar(slide, 0.75, 4.19, 5.35, 0.44, 16384, 65536, GREEN, "Normalize · 16 lanes", "16,384")
metric(slide, 0.92, 5.25, 2.2, "32,256", "overlapped cycles", TEAL)
metric(slide, 3.65, 5.25, 2.2, "66,048", "temporal span", RED)

label_box(slide, 6.75, 1.50, 5.55, 1.08, "1 · 16-result/cycle EXP2 pipeline", "fast exp2 + explicit latency / initiation interval", PALE_RED, RED, 15, 10)
label_box(slide, 6.75, 2.85, 5.55, 1.08, "2 · 专用vector reduction datapath", "MAX/SUM使用ALU与cross-lane tree，不占用EXP资源", PALE_TEAL, TEAL, 15, 10)
label_box(slide, 6.75, 4.20, 5.55, 1.08, "3 · Streaming normalize → DMA", "计算目标约16.9k后，17.6k DMA路径将成为下一瓶颈", PALE_AMBER, AMBER, 15, 10)
add_text(slide, 7.18, 5.76, 4.72, 0.38, "预计 modeled compute：66.1k → 16.9k cycles", 17, GREEN, True, PP_ALIGN.CENTER, font=EN_FONT)
takeaway(slide, "目标不是复制GPU SM，而是采用GPU-like分层并行构建平衡的Softmax专用数据通路。")
source(slide, "16.9k is a model estimate for exp_lanes=16, not a measured result.")
notes(slide, "独立stage统计表明EXP+SUM的active cycles是另外两个阶段的4倍。将exp lanes从4提高到16后，模型预计从66061降到约16909cycle。此时当前17565cycle的DMA ACK路径会成为下一瓶颈，因此优化顺序应是先真实建模EXP latency和II，再做streaming normalize和output DMA。16900是估算，不能当成已经完成的实验结果。")


# 16. Summary
slide = prs.slides.add_slide(blank); set_bg(slide); add_title(slide, "本阶段完成了从“通信驱动”到“行级流水驱动”的架构转变", "Summary", 16)
claims = [
    ("架构", "1 tensor job\n16 Row Engines\n4-context pipeline", TEAL, PALE_TEAL),
    ("数据流", "1 input + 1 output\n0 reduction messages\n64 DMA pairs", BLUE, PALE_BLUE),
    ("正确性", "4,194,304 checked\n0 mismatch\nstrict lifecycle PASS", GREEN, PALE_GREEN),
    ("性能", "66,062 accelerator\n72,409 clean kernel\n27.8× system progress", RED, PALE_RED),
]
for i, (title, body, col, pale) in enumerate(claims):
    x = 0.68 + i * 3.12
    rect(slide, x, 1.62, 2.65, 2.36, pale, GRID, True)
    add_text(slide, x+0.18, 1.91, 2.29, 0.34, title, 17, col, True, PP_ALIGN.CENTER)
    add_text(slide, x+0.18, 2.62, 2.29, 0.86, body, 11.5, INK, False, PP_ALIGN.CENTER, font=EN_FONT)
add_text(slide, 1.08, 4.68, 11.15, 0.44, "当前结论：NoC/DMA已被隐藏，4-lane EXP吞吐决定66k-cycle关键路径。", 17, INK, True, PP_ALIGN.CENTER)
rect(slide, 2.15, 5.55, 9.05, 0.78, INK, INK, True)
add_text(slide, 2.35, 5.76, 8.65, 0.30, "下一阶段验收：modeled ≤ 17k · observed ≤ 18k · golden/lifecycle不变", 14, WHITE, True, PP_ALIGN.CENTER, font=EN_FONT)
takeaway(slide, "汇报重点：我们不仅获得了更少的cycle，也建立了可解释、可验证、可继续优化的模型。", INK, "E4ECE7")
notes(slide, "总结四点：架构上完成单tensor job和16 Row Engines；数据流上只保留输入输出两次HBM访问；测试上全量golden和生命周期通过；性能上accelerator 66062、clean kernel 72409。下一阶段不是盲目DSE，而是实现可信的16-wide EXP pipeline和独立reduction datapath。")


for slide in prs.slides:
    # Keep every shape within slide bounds by construction; add a subtle footer marker.
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.62), Inches(7.30), Inches(0.12), Inches(0.03))
    marker.fill.solid(); marker.fill.fore_color.rgb = rgb(TEAL); marker.line.fill.background()

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(OUT_PATH)
